"""PPTX text extraction using python-pptx with embedded image OCR."""

import io
import os
from models import ExtractionResult, PageContent, ExtractionMetadata
from .base import BaseExtractor
from utils.text_normalizer import normalize_text

try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False


class PPTXExtractor(BaseExtractor):
    """Extract text from PowerPoint presentations slide by slide."""

    async def extract(self, file_path: str, file_name: str) -> ExtractionResult:
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from PIL import Image
        except ImportError as e:
            return self._create_error_result(file_name, "pptx", f"Missing dependency: {e}")

        file_size = os.path.getsize(file_path)

        try:
            prs = Presentation(file_path)
            pages: list[PageContent] = []
            all_text_parts: list[str] = []
            has_images = False

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts: list[str] = []
                slide_title = ""
                image_texts: list[str] = []

                for shape in slide.shapes:
                    # Extract text from text frames
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_texts.append(text)

                    # Get slide title
                    if shape.has_text_frame and shape.shape_id == 0:
                        slide_title = shape.text_frame.text.strip()

                    # Extract text from tables in slides
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            slide_texts.append(" | ".join(cells))

                    # OCR on embedded images
                    if shape.shape_type == 13 and TESSERACT_AVAILABLE:  # Picture type
                        has_images = True
                        try:
                            image_data = shape.image.blob
                            img = Image.open(io.BytesIO(image_data))
                            text = pytesseract.image_to_string(img)
                            if text.strip():
                                image_texts.append(text.strip())
                        except Exception:
                            continue

                # If no explicit title, use first line
                if not slide_title and slide_texts:
                    slide_title = slide_texts[0][:80]

                combined = normalize_text("\n".join(slide_texts))
                image_text_combined = normalize_text("\n".join(image_texts))

                page = PageContent(
                    page_number=slide_num,
                    text=combined,
                    has_images=bool(image_texts),
                    image_text=image_text_combined,
                    slide_title=slide_title,
                )
                pages.append(page)
                all_text_parts.append(f"--- Slide {slide_num}: {slide_title} ---\n{combined}")
                if image_text_combined:
                    all_text_parts.append(image_text_combined)

            extracted = normalize_text("\n\n".join(all_text_parts))
            word_count = len(extracted.split())

            return ExtractionResult(
                success=True,
                file_name=file_name,
                file_type="pptx",
                file_size=file_size,
                total_pages=len(pages),
                extracted_text=extracted,
                pages=pages,
                metadata=ExtractionMetadata(
                    word_count=word_count,
                    char_count=len(extracted),
                    has_images=has_images,
                    extraction_method="python-pptx+OCR",
                ),
            )
        except Exception as e:
            return self._create_error_result(file_name, "pptx", str(e))
